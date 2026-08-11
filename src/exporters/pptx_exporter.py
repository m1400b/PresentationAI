"""
PresentationAI

PowerPoint Exporter
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Cm

from src.models.slide import Slide
from src.renderers.slide_renderer import SlideRenderer


class PPTXExporter:
    """
    Exports a PresentationDocument to PPTX.
    """

    # -------------------------------------------------

    def __init__(self):

        self.slide_renderer = SlideRenderer()

    # -------------------------------------------------

    def export(
        self,
        slides: list[Slide],
        filename: str,
    ):

        prs = Presentation()

        #
        # 16:9
        #

        prs.slide_width = Cm(33.867)

        prs.slide_height = Cm(19.05)

        blank_layout = prs.slide_layouts[6]

        for slide in slides:

            ppt_slide = prs.slides.add_slide(
                blank_layout
            )

            self.slide_renderer.render(

                ppt_slide,

                slide,

            )

        Path(filename).parent.mkdir(

            parents=True,

            exist_ok=True,

        )

        prs.save(filename)

        print(

            f"{len(slides)} slide(s) exported."

        )