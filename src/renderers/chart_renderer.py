"""
PresentationAI

Chart Renderer
"""

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Cm

from src.models.elements.chart_element import ChartElement


class ChartRenderer:

    def render(
        self,
        ppt_slide,
        element: ChartElement,
    ):

        shape = ppt_slide.shapes.add_shape(

            MSO_AUTO_SHAPE_TYPE.RECTANGLE,

            Cm(element.x / 37.8),

            Cm(element.y / 37.8),

            Cm(element.width / 37.8),

            Cm(element.height / 37.8),

        )

        frame = shape.text_frame

        frame.text = f"{element.chart_type.upper()} CHART"