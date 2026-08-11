"""
PresentationAI

Image Renderer
"""

from pathlib import Path

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

from src.models.elements.image_element import ImageElement


class ImageRenderer:
    """
    Renders an ImageElement into PowerPoint.
    """

    PX_TO_CM = 37.7952755906

    # -------------------------------------------------

    @classmethod
    def px(cls, value: float):

        return Cm(value / cls.PX_TO_CM)

    # -------------------------------------------------

    def render(
        self,
        ppt_slide,
        element: ImageElement,
    ):

        #
        # Real image
        #

        if element.path and Path(element.path).exists():

            try:

                ppt_slide.shapes.add_picture(

                    element.path,

                    self.px(element.x),

                    self.px(element.y),

                    width=self.px(element.width),

                    height=self.px(element.height),

                )

                return

            except Exception:

                pass

        #
        # Placeholder
        #

        shape = ppt_slide.shapes.add_shape(

            MSO_AUTO_SHAPE_TYPE.RECTANGLE,

            self.px(element.x),

            self.px(element.y),

            self.px(element.width),

            self.px(element.height),

        )

        text_frame = shape.text_frame

        text_frame.clear()

        text_frame.word_wrap = True

        paragraph = text_frame.paragraphs[0]

        paragraph.alignment = PP_ALIGN.CENTER

        paragraph.font.size = Pt(14)

        if element.caption:

            paragraph.text = (
                "[ IMAGE ]\n\n"
                + element.caption
            )

        else:

            paragraph.text = "[ IMAGE ]"