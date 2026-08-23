"""
PresentationAI

PowerPoint Renderer
"""

from __future__ import annotations

from pathlib import Path


from pptx import Presentation as PptxPresentation

from pptx.util import Inches
from pptx.util import Pt

from pptx.enum.text import PP_ALIGN

from pptx.dml.color import RGBColor


from src.renderers.base_renderer import (
    BaseRenderer,
)

from src.models.presentation import (
    Presentation,
)

from src.models.elements.text_element import (
    TextElement,
)

from src.models.elements.image_element import (
    ImageElement,
)



class PptxRenderer(
    BaseRenderer,
):

    # =================================================
    # Identity
    # =================================================

    @property
    def renderer_name(
        self,
    ) -> str:

        return "pptx"


    # =================================================
    # Render
    # =================================================

    def render(
        self,
        presentation: Presentation,
        output_path: str,
    ) -> str:


        self.validate(
            presentation
        )


        output_path = self.ensure_output(
            output_path
        )


        prs = PptxPresentation()


        #
        # 16:9
        #

        prs.slide_width = Inches(
            13.333
        )

        prs.slide_height = Inches(
            7.5
        )


        for slide in presentation.slides:

            self._render_slide(
                prs,
                slide,
            )


        prs.save(
            output_path
        )


        return output_path



    # =================================================
    # Slide
    # =================================================

    def _render_slide(
        self,
        prs,
        slide,
    ):


        ppt_slide = prs.slides.add_slide(
            prs.slide_layouts[6]
        )


        for element in slide.elements:


            if isinstance(
                element,
                TextElement,
            ):

                self._add_text(
                    ppt_slide,
                    element,
                )


            elif isinstance(
                element,
                ImageElement,
            ):

                self._add_image(
                    ppt_slide,
                    element,
                )



    # =================================================
    # Text
    # =================================================

    def _add_text(
        self,
        ppt_slide,
        element: TextElement,
    ):


        box = ppt_slide.shapes.add_textbox(

            Inches(
                self._x(element.x)
            ),

            Inches(
                self._y(element.y)
            ),

            Inches(
                self._w(element.width)
            ),

            Inches(
                self._h(element.height)
            ),

        )


        frame = box.text_frame

        frame.clear()


        paragraph = frame.paragraphs[0]


        paragraph.alignment = (
            PP_ALIGN.RIGHT
            if element.style.horizontal_alignment == "right"
            else PP_ALIGN.LEFT
        )


        run = paragraph.add_run()


        run.text = element.text


        font = run.font


        font.name = (
            element.style.font_family
        )


        font.size = Pt(
            element.style.font_size
        )


        font.bold = (
            element.style.bold
        )


        font.italic = (
            element.style.italic
        )


        if element.style.color:

            font.color.rgb = (
                self._hex_color(
                    element.style.color
                )
            )



    # =================================================
    # Image
    # =================================================

    def _add_image(
        self,
        ppt_slide,
        element: ImageElement,
    ):


        if not element.path:

            return


        if not Path(
            element.path
        ).exists():

            return


        ppt_slide.shapes.add_picture(

            element.path,

            Inches(
                self._x(element.x)
            ),

            Inches(
                self._y(element.y)
            ),

            Inches(
                self._w(element.width)
            ),

            Inches(
                self._h(element.height)
            ),

        )



    # =================================================
    # Coordinate Mapper
    # =================================================


    def _x(self,x):

        return x * 13.333 / 24


    def _y(self,y):

        return y * 7.5 / 10


    def _w(self,w):

        return w * 13.333 / 24


    def _h(self,h):

        return h * 7.5 / 10



    # =================================================
    # Color
    # =================================================


    def _hex_color(
        self,
        value:str,
    ):

        value = (
            value
            .replace("#","")
        )

        return RGBColor(

            int(value[0:2],16),

            int(value[2:4],16),

            int(value[4:6],16),

        )